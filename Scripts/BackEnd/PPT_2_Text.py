from pptx import Presentation
import os
import zipfile
import tempfile

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def process_pptx(zip_path):
    with tempfile.TemporaryDirectory() as tmpdirname:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(tmpdirname)

        txt_files = []
        for root, _, files in os.walk(tmpdirname):
            for file in files:
                if file.endswith(".pptx"):
                    pptx_path = os.path.join(root, file)
                    text = extract_text_from_pptx(pptx_path)
                    txt_filename = f"{os.path.splitext(file)[0]}.txt"
                    txt_path = os.path.join(root, txt_filename)
                    with open(txt_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                    txt_files.append(txt_path)

        result_zip_path = os.path.join(tmpdirname, 'result.zip')
        with zipfile.ZipFile(result_zip_path, 'w') as zipf:
            for txt_file in txt_files:
                zipf.write(txt_file, os.path.basename(txt_file))

    return result_zip_path
